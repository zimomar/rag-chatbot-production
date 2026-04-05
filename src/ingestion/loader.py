"""
Chargement de documents PDF, Markdown, DOCX et PPTX.

Ce module extrait le texte brut des fichiers uploadés tout en préservant
les métadonnées utiles (source, nombre de pages, date).
"""

import logging
from dataclasses import dataclass, field
from datetime import UTC, datetime
from io import BytesIO
from pathlib import Path
from typing import BinaryIO

import fitz  # PyMuPDF
import markdown
from docx import Document as DocxDocument
from pptx import Presentation

logger = logging.getLogger(__name__)


@dataclass
class Document:
    """
    Représente un document chargé avec son contenu et métadonnées.

    Attributes:
        content: Texte brut extrait du document
        metadata: Informations sur le document (source, pages, date, etc.)
    """

    content: str
    metadata: dict = field(default_factory=dict)

    def __post_init__(self) -> None:
        """Ajoute un timestamp si non présent."""
        if "created_at" not in self.metadata:
            self.metadata["created_at"] = datetime.now(UTC).isoformat()

    @property
    def source(self) -> str:
        """Nom du fichier source."""
        return str(self.metadata.get("source", "unknown"))

    @property
    def num_pages(self) -> int:
        """Nombre de pages (pour PDF)."""
        return int(self.metadata.get("num_pages", 1))

    def __len__(self) -> int:
        """Longueur du contenu en caractères."""
        return len(self.content)


class DocumentLoader:
    """
    Charge et extrait le texte de documents PDF et Markdown.

    Exemples d'utilisation:
        >>> loader = DocumentLoader()
        >>> doc = loader.load_file(Path("document.pdf"))
        >>> print(doc.content[:100])
    """

    SUPPORTED_EXTENSIONS = {".pdf", ".md", ".markdown", ".txt", ".docx", ".pptx"}

    def load_file(self, file_path: Path) -> Document:
        """
        Charge un document depuis un chemin de fichier.

        Args:
            file_path: Chemin vers le fichier à charger

        Returns:
            Document avec le contenu extrait et les métadonnées

        Raises:
            FileNotFoundError: Si le fichier n'existe pas
            ValueError: Si le format n'est pas supporté
        """
        if not file_path.exists():
            raise FileNotFoundError(f"Fichier non trouvé: {file_path}")

        extension = file_path.suffix.lower()
        if extension not in self.SUPPORTED_EXTENSIONS:
            raise ValueError(
                f"Format non supporté: {extension}. "
                f"Formats acceptés: {', '.join(self.SUPPORTED_EXTENSIONS)}"
            )

        logger.info(f"Chargement du fichier: {file_path.name}")

        with open(file_path, "rb") as f:
            return self._load_from_bytes(f, file_path.name, extension)

    def load_uploaded_file(
        self,
        file_content: BinaryIO | bytes,
        filename: str,
    ) -> Document:
        """
        Charge un document depuis un fichier uploadé (bytes ou file-like).

        Args:
            file_content: Contenu du fichier (bytes ou objet file-like)
            filename: Nom du fichier original

        Returns:
            Document avec le contenu extrait et les métadonnées

        Raises:
            ValueError: Si le format n'est pas supporté
        """
        extension = Path(filename).suffix.lower()
        if extension not in self.SUPPORTED_EXTENSIONS:
            raise ValueError(
                f"Format non supporté: {extension}. "
                f"Formats acceptés: {', '.join(self.SUPPORTED_EXTENSIONS)}"
            )

        logger.info(f"Chargement du fichier uploadé: {filename}")

        # Convertit bytes en BytesIO si nécessaire
        if isinstance(file_content, bytes):
            file_content = BytesIO(file_content)

        return self._load_from_bytes(file_content, filename, extension)

    def _load_from_bytes(
        self,
        file_obj: BinaryIO,
        filename: str,
        extension: str,
    ) -> Document:
        """
        Charge un document depuis un objet file-like.

        Args:
            file_obj: Objet file-like contenant les données
            filename: Nom du fichier
            extension: Extension du fichier

        Returns:
            Document avec le contenu extrait
        """
        if extension == ".pdf":
            return self._load_pdf(file_obj, filename)
        elif extension in {".md", ".markdown"}:
            return self._load_markdown(file_obj, filename)
        elif extension == ".docx":
            return self._load_docx(file_obj, filename)
        elif extension == ".pptx":
            return self._load_pptx(file_obj, filename)
        else:  # .txt
            return self._load_text(file_obj, filename)

    def _load_pdf(self, file_obj: BinaryIO, filename: str) -> Document:
        """
        Extrait le texte d'un fichier PDF.

        Utilise PyMuPDF (fitz) pour une extraction rapide et fiable.
        Concatène le texte de toutes les pages avec des séparateurs.

        Args:
            file_obj: Objet file-like contenant le PDF
            filename: Nom du fichier

        Returns:
            Document avec le texte extrait page par page
        """
        content = file_obj.read()
        pages_text = []
        num_pages = 0

        try:
            with fitz.open(stream=content, filetype="pdf") as pdf_doc:
                num_pages = len(pdf_doc)
                logger.debug(f"PDF {filename}: {num_pages} pages")

                for page_num, page in enumerate(pdf_doc, start=1):
                    text = page.get_text("text")
                    if text.strip():
                        # Ajoute un marqueur de page pour le contexte
                        pages_text.append(f"[Page {page_num}]\n{text.strip()}")

        except Exception as e:
            logger.error(f"Erreur lors du parsing PDF {filename}: {e}")
            raise ValueError(f"Impossible de parser le PDF: {e}") from e

        if not pages_text:
            logger.warning(f"Aucun texte extrait de {filename}")
            raise ValueError(
                f"Aucun texte trouvé dans {filename}. "
                "Le PDF est peut-être un scan (image) sans OCR."
            )

        full_text = "\n\n".join(pages_text)

        return Document(
            content=full_text,
            metadata={
                "source": filename,
                "num_pages": num_pages,
                "file_type": "pdf",
            },
        )

    def _load_markdown(self, file_obj: BinaryIO, filename: str) -> Document:
        """
        Charge et convertit un fichier Markdown en texte brut.

        Préserve la structure (headers, listes) mais supprime le formatage.

        Args:
            file_obj: Objet file-like contenant le Markdown
            filename: Nom du fichier

        Returns:
            Document avec le texte converti
        """
        try:
            content = file_obj.read().decode("utf-8")
        except UnicodeDecodeError:
            # Fallback sur latin-1 si UTF-8 échoue
            file_obj.seek(0)
            content = file_obj.read().decode("latin-1")

        # Convertit MD en HTML puis extrait le texte
        # Cela préserve la structure tout en nettoyant le formatage
        md_converter = markdown.Markdown(extensions=["tables", "fenced_code", "nl2br"])
        html = md_converter.convert(content)

        # Nettoyage basique du HTML pour obtenir du texte brut
        # On garde les sauts de ligne pour la structure
        import re

        text = re.sub(r"<[^>]+>", "", html)
        text = re.sub(r"\n{3,}", "\n\n", text)  # Max 2 sauts de ligne

        return Document(
            content=text.strip(),
            metadata={
                "source": filename,
                "num_pages": 1,
                "file_type": "markdown",
            },
        )

    def _load_text(self, file_obj: BinaryIO, filename: str) -> Document:
        """
        Charge un fichier texte brut.

        Args:
            file_obj: Objet file-like contenant le texte
            filename: Nom du fichier

        Returns:
            Document avec le contenu
        """
        try:
            content = file_obj.read().decode("utf-8")
        except UnicodeDecodeError:
            file_obj.seek(0)
            content = file_obj.read().decode("latin-1")

        return Document(
            content=content.strip(),
            metadata={
                "source": filename,
                "num_pages": 1,
                "file_type": "text",
            },
        )

    def _load_docx(self, file_obj: BinaryIO, filename: str) -> Document:
        """
        Extrait le texte d'un fichier Word (.docx).

        Parcourt les paragraphes et tableaux du document.

        Args:
            file_obj: Objet file-like contenant le DOCX
            filename: Nom du fichier

        Returns:
            Document avec le texte extrait
        """
        try:
            content = file_obj.read()
            docx_doc = DocxDocument(BytesIO(content))
        except Exception as e:
            logger.error(f"Erreur lors du parsing DOCX {filename}: {e}")
            raise ValueError(f"Impossible de parser le DOCX: {e}") from e

        paragraphs = []

        # Extraire les paragraphes
        for para in docx_doc.paragraphs:
            text = para.text.strip()
            if text:
                # Préserve la hiérarchie des headings
                if para.style and para.style.name.startswith("Heading"):
                    level = para.style.name.replace("Heading ", "")
                    try:
                        paragraphs.append(f"{'#' * int(level)} {text}")
                    except ValueError:
                        paragraphs.append(text)
                else:
                    paragraphs.append(text)

        # Extraire le texte des tableaux
        for table in docx_doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    paragraphs.append(row_text)

        if not paragraphs:
            raise ValueError(f"Aucun texte trouvé dans {filename}.")

        full_text = "\n\n".join(paragraphs)

        return Document(
            content=full_text,
            metadata={
                "source": filename,
                "num_pages": 1,
                "file_type": "docx",
            },
        )

    def _load_pptx(self, file_obj: BinaryIO, filename: str) -> Document:
        """
        Extrait le texte d'un fichier PowerPoint (.pptx).

        Parcourt toutes les slides et en extrait le texte.

        Args:
            file_obj: Objet file-like contenant le PPTX
            filename: Nom du fichier

        Returns:
            Document avec le texte extrait slide par slide
        """
        try:
            content = file_obj.read()
            prs = Presentation(BytesIO(content))
        except Exception as e:
            logger.error(f"Erreur lors du parsing PPTX {filename}: {e}")
            raise ValueError(f"Impossible de parser le PPTX: {e}") from e

        slides_text = []
        num_slides = len(prs.slides)

        for slide_num, slide in enumerate(prs.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            texts.append(text)
            if texts:
                slides_text.append(f"[Slide {slide_num}]\n" + "\n".join(texts))

        if not slides_text:
            raise ValueError(f"Aucun texte trouvé dans {filename}.")

        full_text = "\n\n".join(slides_text)

        return Document(
            content=full_text,
            metadata={
                "source": filename,
                "num_pages": num_slides,
                "file_type": "pptx",
            },
        )
